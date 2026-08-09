from __future__ import annotations

import csv
import importlib.util
import io
import re
import zipfile
from abc import ABC, abstractmethod
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from docx import Document as DocxDocument
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.pipeline.documents.domain.analysis import (
    DocumentAnalyzer,
    ExtractionStrategy,
    unlock_pdf_with_empty_password,
)
from app.pipeline.documents.domain.models import LogicalDocumentFactory
from app.pipeline.documents.domain.parsed import (
    ParsedDocument,
    ParsedElement,
    ParsedImageMetadata,
    ParsedPage,
    ParsedSection,
    ParsedTable,
)
from app.pipeline.documents.extraction.parsing.native_pdf import extract_native_pdf_structure
from app.pipeline.shared.errors import AppError
from app.pipeline.shared.markdown import (
    MARKDOWN_REPRESENTATION_VERSION,
    render_parsed_document_markdown,
)
from app.pipeline.shared.table_text import render_markdown_table_text, render_table_text
from app.pipeline.shared.text_utils import (
    detect_language,
    detect_prompt_injection_like_content,
    normalize_text,
)


class BaseParser(ABC):
    parser_name = "base"
    parser_version = "1.0"
    supported_extensions: tuple[str, ...] = ()

    def supports(self, extension: str) -> bool:
        return extension.lower() in self.supported_extensions

    def validate(self, _content: bytes) -> None:
        return None

    def extract_metadata(self, parsed: ParsedDocument) -> dict[str, str | int | float | bool]:
        return {
            "page_count": len(parsed.pages),
            "word_count": len(parsed.text.split()),
            "table_count": len(parsed.tables),
            "image_count": len(parsed.images_metadata),
            "parser_name": self.parser_name,
            "parser_version": self.parser_version,
            "detected_language": parsed.detected_language,
            "ocr_used": parsed.ocr_used,
        }

    @abstractmethod
    def parse(self, content: bytes) -> ParsedDocument:
        raise NotImplementedError

    def _build_document(
        self,
        *,
        text: str,
        sections: list[ParsedSection],
        pages: list[ParsedPage] | None = None,
        tables: list[ParsedTable] | None = None,
        images_metadata: list[ParsedImageMetadata] | None = None,
        warnings: list[str] | None = None,
        confidence: float | None = None,
        ocr_used: bool = False,
        metadata: dict[str, str | int | float | bool] | None = None,
    ) -> ParsedDocument:
        normalized_text = normalize_text(text)
        all_warnings = list(warnings or [])
        all_warnings.extend(detect_prompt_injection_like_content(normalized_text))
        document = ParsedDocument(
            text=normalized_text,
            pages=list(pages or []),
            sections=sections,
            tables=list(tables or []),
            images_metadata=list(images_metadata or []),
            document_metadata=dict(metadata or {}),
            warnings=all_warnings,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
            confidence=confidence,
            ocr_used=ocr_used,
            detected_language=detect_language(normalized_text),
        )
        document.document_metadata.update(self.extract_metadata(document))
        document.document_metadata["content_format"] = "markdown"
        document.document_metadata["content_representation_version"] = (
            MARKDOWN_REPRESENTATION_VERSION
        )
        document.content_markdown = render_parsed_document_markdown(document)
        document.logical_document = LogicalDocumentFactory.from_parsed_parts(
            text=document.text,
            title=str(document.document_metadata.get("title") or self.parser_name.upper()),
            pages=document.pages,
            sections=document.sections,
            tables=document.tables,
            images_metadata=document.images_metadata,
            metadata=document.document_metadata,
            parser_name=document.parser_name,
            parser_version=document.parser_version,
            warnings=document.warnings,
            confidence=document.confidence,
            ocr_used=document.ocr_used,
            detected_language=document.detected_language,
        )
        return document


class TxtParser(BaseParser):
    parser_name = "txt"
    supported_extensions = ("txt",)

    def parse(self, content: bytes) -> ParsedDocument:
        text, decode_warning = _decode_text_lossless(content, preferred="utf-8")
        text = normalize_text(text)
        warnings = [decode_warning] if decode_warning else []
        return self._build_document(
            text=text,
            sections=[ParsedSection(text=text, page_number=1, title="Text")],
            pages=[ParsedPage(page_number=1, text=text)],
            warnings=warnings,
        )


class DocxParser(BaseParser):
    parser_name = "docx"
    parser_version = "2.1"
    supported_extensions = ("docx",)

    def validate(self, content: bytes) -> None:
        if not zipfile.is_zipfile(io.BytesIO(content)):
            raise AppError("corrupted_file", "DOCX khong hop le hoac da bi hong.", status_code=422)

    def parse(self, content: bytes) -> ParsedDocument:
        self.validate(content)
        doc = DocxDocument(io.BytesIO(content))
        elements: list[ParsedElement] = []
        logical_parts: list[str] = []
        sections: list[ParsedSection] = []
        current_title = "DOCX"
        current_level = 1
        current_section_parts: list[str] = []
        current_section_block_ids: list[str] = []
        tables: list[ParsedTable] = []
        table_index = 0

        def flush_section() -> None:
            if not current_section_block_ids:
                return
            sections.append(
                ParsedSection(
                    text=normalize_text("\n\n".join(current_section_parts)),
                    page_number=1,
                    title=current_title,
                    level=current_level,
                    block_ids=list(current_section_block_ids),
                )
            )
            current_section_parts.clear()
            current_section_block_ids.clear()

        for element_index, item in enumerate(doc.iter_inner_content(), start=1):
            if isinstance(item, DocxParagraph):
                paragraph_text = normalize_text(item.text)
                if not paragraph_text:
                    continue
                style_name = str(getattr(item.style, "name", "") or "")
                heading_match = re.fullmatch(r"Heading\s+([1-6])", style_name, re.IGNORECASE)
                if heading_match:
                    flush_section()
                    current_title = paragraph_text
                    current_level = int(heading_match.group(1))
                    block_type = "heading"
                else:
                    current_section_parts.append(paragraph_text)
                    block_type = "paragraph"
                element_id = f"docx-element-{element_index}"
                elements.append(
                    ParsedElement(
                        element_id=element_id,
                        block_type=block_type,
                        text=paragraph_text,
                        page_number=1,
                    )
                )
                current_section_block_ids.append(element_id)
                logical_parts.append(paragraph_text)
                continue
            if not isinstance(item, DocxTable):
                continue
            rows = [[normalize_text(cell.text) for cell in row.cells] for row in item.rows]
            if not any(any(cell for cell in row) for row in rows):
                continue
            table_text = _render_table(rows)
            table = ParsedTable(
                table_id=f"docx-table-{table_index}",
                location=f"document:page:1:table:{table_index}",
                rows=rows,
                columns=max((len(row) for row in rows), default=0),
                header=rows[0] if rows else [],
            )
            table_index += 1
            tables.append(table)
            current_section_parts.append(table_text)
            logical_parts.append(table_text)
            elements.append(_table_element(table, page_number=1))
            current_section_block_ids.append(table.table_id)
        flush_section()
        text = normalize_text("\n\n".join(logical_parts))
        return self._build_document(
            text=text,
            sections=sections or [ParsedSection(text=text, page_number=1, title="DOCX")],
            pages=[ParsedPage(page_number=1, text=text, elements=elements)],
            tables=tables,
        )


class PdfParser(BaseParser):
    parser_name = "pdf"
    parser_version = "2.0"
    supported_extensions = ("pdf",)

    def parse(self, content: bytes) -> ParsedDocument:
        analysis = DocumentAnalyzer().analyze("document.pdf", content)
        if analysis.encrypted:
            raise AppError(
                "password_protected_file", "PDF dang duoc bao ve bang mat khau.", status_code=422
            )
        if analysis.corrupted:
            raise AppError("corrupted_file", "PDF khong hop le hoac da bi hong.", status_code=422)
        if analysis.extraction_strategy == ExtractionStrategy.OCR.value:
            raise AppError(
                "ocr_not_available",
                "PDF can OCR nhung OCR chua duoc cau hinh cho moi truong hien tai.",
                status_code=503,
            )
        reader = PdfReader(io.BytesIO(content))
        if not unlock_pdf_with_empty_password(reader):
            raise AppError(
                "password_protected_file",
                "PDF dang duoc bao ve bang mat khau.",
                status_code=422,
            )
        sections: list[ParsedSection] = []
        pages: list[ParsedPage] = []
        all_text: list[str] = []
        warnings: list[str] = []
        fallback_page_texts: dict[int, str] = {}
        for idx, page in enumerate(reader.pages, start=1):
            page_text = normalize_text(page.extract_text() or "")
            fallback_page_texts[idx] = page_text
            if page_text:
                sections.append(ParsedSection(text=page_text, page_number=idx, title=f"Page {idx}"))
                pages.append(ParsedPage(page_number=idx, text=page_text))
                all_text.append(page_text)
        if not all_text:
            warnings.append("ocr_required_not_production_ready")
            raise AppError(
                "parser_failure",
                "Khong the trich xuat noi dung tu PDF text-native.",
                status_code=422,
            )
        native = extract_native_pdf_structure(
            content,
            fallback_page_texts=fallback_page_texts,
        )
        base_metadata = {
            "document_analysis": analysis.to_dict(),
            "pdf_type": analysis.pdf_type or "unknown",
            "extraction_strategy": analysis.extraction_strategy,
            "estimated_ocr_required": analysis.estimated_ocr_required,
            "analysis_confidence": analysis.confidence,
        }
        if native is not None and native.text.strip():
            return self._build_document(
                text=native.text,
                sections=native.sections,
                pages=native.pages,
                tables=native.tables,
                images_metadata=native.images_metadata,
                warnings=[*warnings, *native.warnings],
                metadata={**base_metadata, **native.metadata},
            )
        return self._build_document(
            text="\n\n".join(all_text),
            sections=sections,
            pages=pages,
            warnings=warnings,
            metadata={
                **base_metadata,
                "native_pdf_extraction": {
                    "provider": "pypdf",
                    "version": self.parser_version,
                    "mode": "text_only_fallback",
                    "text_source": "pypdf",
                    "page_count": len(pages),
                    "line_count": 0,
                    "element_count": 0,
                    "heading_count": 0,
                    "table_candidate_count": 0,
                    "table_count": 0,
                    "image_count": 0,
                    "model_fallback_status": "not_configured",
                },
            },
        )


class PptxParser(BaseParser):
    parser_name = "pptx"
    parser_version = "2.1"
    supported_extensions = ("pptx",)

    def validate(self, content: bytes) -> None:
        if not zipfile.is_zipfile(io.BytesIO(content)):
            raise AppError("corrupted_file", "PPTX khong hop le hoac da bi hong.", status_code=422)

    def parse(self, content: bytes) -> ParsedDocument:
        self.validate(content)
        presentation = Presentation(io.BytesIO(content))
        sections: list[ParsedSection] = []
        pages: list[ParsedPage] = []
        tables: list[ParsedTable] = []
        images_metadata: list[ParsedImageMetadata] = []
        slide_texts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            elements: list[ParsedElement] = []
            title_shape = slide.shapes.title
            slide_title = normalize_text(title_shape.text) if title_shape is not None else ""
            for shape_index, shape in enumerate(slide.shapes):
                if getattr(shape, "has_table", False):
                    rows = [
                        [normalize_text(cell.text) for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if not any(any(cell for cell in row) for row in rows):
                        continue
                    table = ParsedTable(
                        table_id=f"pptx-table-{slide_index}-{shape_index}",
                        location=f"slide:{slide_index}:page:{slide_index}:table:{shape_index}",
                        rows=rows,
                        columns=max((len(row) for row in rows), default=0),
                        header=rows[0] if rows else [],
                    )
                    tables.append(table)
                    table_text = _render_table(rows)
                    lines.append(table_text)
                    elements.append(_table_element(table, page_number=slide_index))
                elif getattr(shape, "has_text_frame", False):
                    text = normalize_text(shape.text or "")
                    if text:
                        lines.append(text)
                        elements.append(
                            ParsedElement(
                                element_id=f"pptx-element-{slide_index}-{shape_index}",
                                block_type="heading" if shape is title_shape else "paragraph",
                                text=text,
                                page_number=slide_index,
                            )
                        )
                if shape.shape_type == 13:
                    images_metadata.append(
                        ParsedImageMetadata(
                            location=f"slide:{slide_index}:image:{shape_index}", name=shape.name
                        )
                    )
            slide_text = normalize_text("\n".join(lines))
            if slide_text or elements:
                sections.append(
                    ParsedSection(
                        text=slide_text,
                        page_number=slide_index,
                        title=slide_title or f"Slide {slide_index}",
                        block_ids=[element.element_id for element in elements],
                    )
                )
                pages.append(
                    ParsedPage(page_number=slide_index, text=slide_text, elements=elements)
                )
                slide_texts.append(slide_text)
        if not slide_texts and not tables:
            raise AppError(
                "parser_failure", "Khong the trich xuat noi dung tu PPTX.", status_code=422
            )
        return self._build_document(
            text="\n\n".join(slide_texts),
            sections=sections or [ParsedSection(text="", page_number=1, title="PPTX")],
            pages=pages,
            tables=tables,
            images_metadata=images_metadata,
        )


class XlsxParser(BaseParser):
    parser_name = "xlsx"
    parser_version = "2.1"
    supported_extensions = ("xlsx",)

    def validate(self, content: bytes) -> None:
        if not zipfile.is_zipfile(io.BytesIO(content)):
            raise AppError("corrupted_file", "XLSX khong hop le hoac da bi hong.", status_code=422)

    def parse(self, content: bytes) -> ParsedDocument:
        self.validate(content)
        values_workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
        formula_workbook = load_workbook(io.BytesIO(content), data_only=False, read_only=True)
        sections: list[ParsedSection] = []
        pages: list[ParsedPage] = []
        tables: list[ParsedTable] = []
        full_text: list[str] = []
        try:
            for sheet_index, sheet_name in enumerate(values_workbook.sheetnames, start=1):
                sheet = values_workbook[sheet_name]
                formula_sheet = formula_workbook[sheet_name]
                all_rows: list[list[str]] = []
                value_rows = list(sheet.iter_rows(values_only=True))
                formula_rows = list(formula_sheet.iter_rows(values_only=True))
                for value_row, formula_row in zip(value_rows, formula_rows, strict=False):
                    width = max(len(value_row), len(formula_row))
                    row_values: list[str] = []
                    for index in range(width):
                        value = value_row[index] if index < len(value_row) else None
                        formula_value = formula_row[index] if index < len(formula_row) else None
                        row_values.append(normalize_text(_xlsx_cell_text(value, formula_value)))
                    all_rows.append(row_values)
                while all_rows and not any(all_rows[0]):
                    all_rows.pop(0)
                while all_rows and not any(all_rows[-1]):
                    all_rows.pop()
                rows = all_rows
                sheet_text = _render_table(rows)
                page = ParsedPage(page_number=sheet_index, text=sheet_text)
                if rows:
                    table = ParsedTable(
                        table_id=f"xlsx-table-{sheet_index}",
                        location=f"sheet:{sheet_name}:page:{sheet_index}",
                        rows=rows,
                        columns=max((len(row) for row in rows), default=0),
                        header=rows[0] if rows else [],
                    )
                    tables.append(table)
                    page.elements.append(_table_element(table, page_number=sheet_index))
                    full_text.append(f"{sheet_name}\n{sheet_text}")
                sections.append(
                    ParsedSection(
                        text=sheet_text,
                        page_number=sheet_index,
                        title=sheet_name,
                        block_ids=[element.element_id for element in page.elements],
                    )
                )
                pages.append(page)
        finally:
            values_workbook.close()
            formula_workbook.close()
        return self._build_document(
            text="\n\n".join(full_text),
            sections=sections,
            pages=pages,
            tables=tables,
            metadata={"sheet_count": len(values_workbook.sheetnames)},
        )


class CsvParser(BaseParser):
    parser_name = "csv"
    parser_version = "2.1"
    supported_extensions = ("csv",)

    def parse(self, content: bytes) -> ParsedDocument:
        decoded, decode_warning = _decode_text_lossless(content, preferred="utf-8-sig")
        warnings = [decode_warning] if decode_warning else []
        try:
            dialect = csv.Sniffer().sniff(decoded[:8192], delimiters=",;\t")
        except csv.Error:
            dialect = csv.excel
        reader = csv.reader(io.StringIO(decoded, newline=""), dialect)
        rows = [[normalize_text(cell) for cell in row] for row in reader]
        filtered = [row for row in rows if any(row)]
        text = _render_table(filtered)
        tables = [
            ParsedTable(
                table_id="csv-table-1",
                location="sheet:csv",
                rows=filtered,
                columns=max((len(row) for row in filtered), default=0),
                header=filtered[0] if filtered else [],
            )
        ]
        return self._build_document(
            text=text,
            sections=[
                ParsedSection(
                    text=text,
                    page_number=1,
                    title="CSV",
                    block_ids=["csv-table-1"] if filtered else [],
                )
            ],
            pages=[
                ParsedPage(
                    page_number=1,
                    text=text,
                    elements=[_table_element(tables[0], page_number=1)] if filtered else [],
                )
            ],
            tables=tables if filtered else [],
            warnings=warnings,
        )


class MarkdownParser(BaseParser):
    parser_name = "markdown"
    parser_version = "2.1"
    supported_extensions = ("md", "markdown")

    def parse(self, content: bytes) -> ParsedDocument:
        decoded, decode_warning = _decode_text_lossless(content, preferred="utf-8")
        warnings = [decode_warning] if decode_warning else []
        sections: list[ParsedSection] = []
        current_title = "Markdown"
        current_level = 1
        current_lines: list[str] = []
        current_block_ids: list[str] = []
        elements: list[ParsedElement] = []
        heading_index = 0
        paragraph_index = 0

        def flush_section() -> None:
            if not current_block_ids:
                return
            sections.append(
                ParsedSection(
                    text=normalize_text("\n".join(current_lines)),
                    page_number=1,
                    title=current_title,
                    level=current_level,
                    block_ids=list(current_block_ids),
                )
            )
            current_lines.clear()
            current_block_ids.clear()

        for raw_line in decoded.splitlines():
            line = normalize_text(raw_line)
            heading = re.fullmatch(r"(#{1,6})[ \t]+(.+?)(?:[ \t]+#*)?", line)
            if heading:
                flush_section()
                current_title = heading.group(2).strip()
                current_level = len(heading.group(1))
                heading_index += 1
                element_id = f"markdown-heading-{heading_index}"
                elements.append(
                    ParsedElement(
                        element_id=element_id,
                        block_type="heading",
                        text=current_title,
                        page_number=1,
                    )
                )
                current_block_ids.append(element_id)
            elif line:
                current_lines.append(line)
                paragraph_index += 1
                element_id = f"markdown-paragraph-{paragraph_index}"
                elements.append(
                    ParsedElement(
                        element_id=element_id,
                        block_type="paragraph",
                        text=line,
                        page_number=1,
                    )
                )
                current_block_ids.append(element_id)
        flush_section()
        text = "\n\n".join(element.text for element in elements)
        return self._build_document(
            text=text,
            sections=sections or [ParsedSection(text="", page_number=1, title="Markdown")],
            pages=[ParsedPage(page_number=1, text=text, elements=elements)],
            warnings=warnings,
        )


class HtmlParser(BaseParser):
    parser_name = "html"
    parser_version = "1.2"
    supported_extensions = ("html", "htm")

    def parse(self, content: bytes) -> ParsedDocument:
        decoded, decode_warning = _decode_text_lossless(content, preferred="utf-8")
        soup = BeautifulSoup(decoded, "html.parser")
        warnings: list[str] = [decode_warning] if decode_warning else []
        if soup.find("script") or soup.find("style"):
            warnings.append("active_content_reference")
        for active in soup.find_all(["script", "style"]):
            active.decompose()

        tables: list[ParsedTable] = []
        elements: list[ParsedElement] = []
        element_index = 0
        table_index = 0

        def add_element(block_type: str, value: str) -> None:
            nonlocal element_index
            text_value = normalize_text(value)
            if not text_value:
                return
            element_index += 1
            elements.append(
                ParsedElement(
                    element_id=f"html-{block_type}-{element_index}",
                    block_type=block_type,
                    text=text_value,
                    page_number=1,
                )
            )

        def add_table(table: Tag) -> None:
            nonlocal table_index
            table_index += 1
            rows: list[list[str]] = []
            for row in table.find_all("tr"):
                cells = [
                    normalize_text(cell.get_text(" ", strip=True))
                    for cell in row.find_all(["th", "td"], recursive=False)
                ]
                if any(cells):
                    rows.append(cells)
            if rows:
                parsed_table = ParsedTable(
                    table_id=f"html-table-{table_index}",
                    location=f"html:table:{table_index}:page:1",
                    rows=rows,
                    columns=max((len(row) for row in rows), default=0),
                    header=rows[0] if rows else [],
                )
                tables.append(parsed_table)
                elements.append(_table_element(parsed_table, page_number=1))

        recognized = {
            "p",
            "li",
            "blockquote",
            "pre",
            "h1",
            "h2",
            "h3",
            "h4",
            "h5",
            "h6",
            "table",
        }

        def visit(node: object) -> None:
            if isinstance(node, NavigableString):
                add_element("paragraph", str(node))
                return
            if not isinstance(node, Tag):
                return
            name = node.name.lower()
            if name == "table":
                add_table(node)
                return
            if name in {"h1", "h2", "h3", "h4", "h5", "h6"}:
                add_element("heading", node.get_text(" ", strip=True))
                return
            if name in {"p", "li", "blockquote", "pre"}:
                add_element("paragraph", node.get_text(" ", strip=True))
                return
            if not node.find(list(recognized)):
                add_element("paragraph", node.get_text(" ", strip=True))
                return
            for child in node.children:
                visit(child)

        title = (
            normalize_text(soup.title.get_text())
            if soup.title and soup.title.get_text()
            else "HTML"
        )
        root = soup.body or soup
        for child in root.children:
            visit(child)

        text_value = "\n\n".join(element.text for element in elements)
        sections: list[ParsedSection] = []
        current_title = title
        current_elements: list[ParsedElement] = []

        def flush_section() -> None:
            if not current_elements:
                return
            sections.append(
                ParsedSection(
                    text="\n\n".join(element.text for element in current_elements),
                    page_number=1,
                    title=current_title,
                    block_ids=[element.element_id for element in current_elements],
                )
            )

        for element in elements:
            if element.block_type == "heading":
                flush_section()
                current_elements = []
                current_title = element.text
            current_elements.append(element)
        flush_section()
        return self._build_document(
            text=text_value,
            sections=sections or [ParsedSection(text=text_value, page_number=1, title=title)],
            pages=[ParsedPage(page_number=1, text=text_value, elements=elements)],
            tables=tables,
            warnings=warnings,
        )


class ParserRegistry:
    def __init__(
        self,
        *,
        parsers: list[BaseParser] | tuple[BaseParser, ...] | None = None,
        ocr_enabled: bool = False,
    ) -> None:
        self.ocr_enabled = ocr_enabled
        parser_values = (
            list(parsers)
            if parsers is not None
            else [
                TxtParser(),
                DocxParser(),
                PdfParser(),
                PptxParser(),
                XlsxParser(),
                CsvParser(),
                MarkdownParser(),
                HtmlParser(),
            ]
        )
        self._parsers_by_extension: dict[str, BaseParser] = {}
        for parser in parser_values:
            for extension in parser.supported_extensions:
                normalized = extension.lower().lstrip(".")
                if normalized in self._parsers_by_extension:
                    raise ValueError(f"Duplicate parser registration for extension: {normalized}")
                self._parsers_by_extension[normalized] = parser

    def get_parser(self, filename: str) -> BaseParser:
        extension = Path(filename).suffix.lower().lstrip(".")
        parser = self._parsers_by_extension.get(extension)
        if parser is not None:
            return parser
        raise AppError("unsupported_file", "Dinh dang file chua duoc ho tro.", status_code=415)

    def supports(self, filename: str) -> bool:
        extension = Path(filename).suffix.lower().lstrip(".")
        return extension in self._parsers_by_extension

    @property
    def supported_extensions(self) -> tuple[str, ...]:
        return tuple(sorted(self._parsers_by_extension))

    def check_health(self) -> dict[str, bool]:
        ocr_ready = not self.ocr_enabled or all(
            importlib.util.find_spec(module_name) is not None
            for module_name in ("fitz", "PIL", "paddle", "paddleocr")
        )
        return {
            "txt": True,
            "docx": True,
            "pdf": True,
            "pptx": True,
            "xlsx": True,
            "csv": True,
            "markdown": True,
            "html": True,
            "ocr": ocr_ready,
        }


def _render_table(rows: list[list[str]]) -> str:
    return render_table_text(rows)


def _render_markdown_table(rows: list[list[str]]) -> str:
    return render_markdown_table_text(rows)


def _decode_text_lossless(content: bytes, *, preferred: str) -> tuple[str, str | None]:
    for encoding in (preferred, "utf-8-sig", "utf-8", "cp1258", "cp1252", "latin-1"):
        try:
            return content.decode(
                encoding
            ), None if encoding == preferred else f"decoded_with_{encoding}"
        except UnicodeDecodeError:
            continue
    return content.decode(preferred, errors="replace"), "decode_replacement_characters"


def _xlsx_cell_text(value: object, formula_value: object) -> str:
    if value is not None:
        return str(value)
    if isinstance(formula_value, str) and formula_value.startswith("="):
        return formula_value
    if formula_value is not None:
        return str(formula_value)
    return ""


def _table_element(table: ParsedTable, *, page_number: int) -> ParsedElement:
    metadata = {
        "location": table.location,
        "columns": table.columns,
        "header": list(table.header),
        "canonical_text": _render_table(table.rows),
        "markdown_text": _render_markdown_table(table.rows),
        "cells": [dict(cell) for cell in list(getattr(table, "cells", []) or [])],
        **dict(getattr(table, "metadata", {}) or {}),
    }
    return ParsedElement(
        element_id=table.table_id,
        block_type="table",
        text=_render_table(table.rows),
        page_number=page_number,
        metadata=metadata,
        bbox=table.bbox,
        confidence=table.confidence,
    )
